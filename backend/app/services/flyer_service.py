"""Branded flyer generation (PPTX + PDF) for listing marketing.

Ported from gor-marketing's flyer_builder.py and flyer_pdf.py, adapted for
multi-tenant operation. All branding (logo, colors, brokerage info) is
parameterized via BrandingConfig, sourced from BrandProfile.settings or
Tenant.settings JSONB.

Layout:
- Header with logo and headline ("Just Listed", "Price Reduced", etc.)
- Full-width hero photo with address overlay
- Optional secondary photo pair
- Price/specs bar
- Two-column feature checklist
- AI-generated body copy (italic, centered)
- Agent contact info with QR code and logo
- Branded three-column footer
"""

import io
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

import structlog
from PIL import Image

logger = structlog.get_logger()


@dataclass
class BrandingConfig:
    """Visual branding parameters for flyer generation.

    Populate from BrandProfile.settings["flyer"] or Tenant.settings["flyer"].
    All fields have sensible defaults so flyers render even with minimal config.
    """

    brokerage_name: str = "Your Brokerage"
    brokerage_address: str = ""
    brokerage_phone: str = ""
    brokerage_website: str = ""
    tagline: str = ""
    logo_path: str | None = None
    accent_color_hex: str = "#CC0000"
    headline: str = "Just Listed"
    qr_base_url: str = ""

    @classmethod
    def from_settings(cls, settings: dict, **overrides) -> "BrandingConfig":
        """Build from BrandProfile.settings or Tenant.settings JSONB.

        Expected structure:
            settings = {
                "flyer": {
                    "brokerage_name": "...",
                    "brokerage_address": "...",
                    "brokerage_phone": "...",
                    "brokerage_website": "...",
                    "tagline": "...",
                    "logo_path": "/path/or/url",
                    "accent_color": "#CC0000",
                    "qr_base_url": "https://..."
                }
            }
        """
        flyer = settings.get("flyer", {})
        config = cls(
            brokerage_name=flyer.get("brokerage_name", cls.brokerage_name),
            brokerage_address=flyer.get("brokerage_address", cls.brokerage_address),
            brokerage_phone=flyer.get("brokerage_phone", cls.brokerage_phone),
            brokerage_website=flyer.get("brokerage_website", cls.brokerage_website),
            tagline=flyer.get("tagline", cls.tagline),
            logo_path=flyer.get("logo_path"),
            accent_color_hex=flyer.get("accent_color", cls.accent_color_hex),
            qr_base_url=flyer.get("qr_base_url", cls.qr_base_url),
        )
        for k, v in overrides.items():
            if hasattr(config, k) and v is not None:
                setattr(config, k, v)
        return config

    @property
    def accent_rgb(self) -> tuple[int, int, int]:
        """Parse hex color to (R, G, B) tuple."""
        h = self.accent_color_hex.lstrip("#")
        if len(h) == 6:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        return (0xCC, 0x00, 0x00)

    @property
    def has_logo(self) -> bool:
        return bool(self.logo_path) and Path(self.logo_path).exists()


# ── Shared Helpers ────────────────────────────────────────────────────

def _extract_body_copy(flyer_text: str) -> str:
    """Extract prose body from AI-generated flyer copy, capped at ~80 words."""
    clean_lines = []
    for line in flyer_text.split("\n"):
        s = line.strip()
        if s.startswith("#") or s.startswith("---") or not s:
            continue
        s = s.replace("**", "").replace("*", "")
        if s.startswith("- ") or s.startswith("+ "):
            s = s[2:]
        clean_lines.append(s)

    body = ""
    for line in clean_lines:
        if len(line) > len(body) and len(line) > 60:
            body = line

    if not body:
        body = " ".join(clean_lines[:3])

    words = body.split()
    if len(words) > 80:
        truncated = " ".join(words[:80])
        for end_char in [". ", "! ", "? "]:
            last_end = truncated.rfind(end_char)
            if last_end > len(truncated) // 3:
                return truncated[:last_end + 1]
        for char in [".", ","]:
            last = truncated.rfind(char)
            if last > len(truncated) // 3:
                return truncated[:last + 1]
        body = truncated + "..."

    return body


def _generate_qr(data: str) -> str:
    """Generate a QR code PNG and return the temp file path."""
    import qrcode

    qr = qrcode.QRCode(
        version=1, box_size=10, border=2,
        error_correction=qrcode.constants.ERROR_CORRECT_M,
    )
    qr.add_data(data)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")

    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    img.save(tmp.name)
    return tmp.name


def _build_specs(listing_data: dict) -> list[str]:
    """Build specs list from listing data."""
    specs = []
    price = listing_data.get("price")
    if price:
        specs.append(f"${float(price):,.0f}")
    beds = listing_data.get("bedrooms", "?")
    baths = listing_data.get("bathrooms", "?")
    specs.append(f"{beds} Bed / {baths} Bath")
    sqft = listing_data.get("sqft")
    if sqft:
        specs.append(f"{int(sqft):,} Sq Ft")
    yb = listing_data.get("year_built")
    if yb:
        specs.append(f"Built {yb}")
    return specs


def _build_qr_url(listing_data: dict, branding: BrandingConfig) -> str:
    """Build URL for QR code from branding config or listing address."""
    address = listing_data.get("address_full", "")
    if branding.qr_base_url:
        mls_id = listing_data.get("mls_listing_id", "")
        return f"{branding.qr_base_url.rstrip('/')}/{mls_id}" if mls_id else branding.qr_base_url
    return f"https://www.google.com/maps/search/{address.replace(' ', '+')}"


def _sanitize_text(text: str) -> str:
    """Replace Unicode characters unsupported by PDF core fonts (Latin-1 only)."""
    replacements = {
        "\u2014": "--", "\u2013": "-", "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"', "\u2026": "...", "\u2022": "-",
        "\u25a0": "\u00bb", "\u2032": "'", "\u2033": '"', "\u00a0": " ",
    }
    for char, replacement in replacements.items():
        text = text.replace(char, replacement)
    return text.encode("latin-1", errors="replace").decode("latin-1")


# ── PPTX Builder ──────────────────────────────────────────────────────

def build_flyer_pptx(
    listing_data: dict,
    flyer_text: str,
    branding: BrandingConfig,
    photo_paths: list[Path] | None = None,
) -> io.BytesIO:
    """Generate a branded PPTX flyer. Returns bytes buffer.

    Args:
        listing_data: Listing data dict (or Listing model attributes).
        flyer_text: AI-generated flyer content.
        branding: Visual branding configuration.
        photo_paths: Local image file paths (hero + secondary).

    Returns:
        BytesIO buffer containing the PPTX file.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    BLACK = RGBColor(0x00, 0x00, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_GREY = RGBColor(0x33, 0x33, 0x33)
    MID_GREY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)
    ACCENT = RGBColor(*branding.accent_rgb)

    SLIDE_WIDTH = Inches(8.5)
    SLIDE_HEIGHT = Inches(11)

    def _textbox(slide, left, top, width, height, text, size=12,
                 bold=False, italic=False, color=BLACK, align=PP_ALIGN.LEFT,
                 font="Arial", valign=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        tf.vertical_anchor = valign
        return box

    def _multiline_textbox(slide, left, top, width, height, lines, size=10,
                           color=BLACK, align=PP_ALIGN.LEFT, font="Arial", bold=False):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        for i, line_text in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line_text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font
            p.alignment = align
            p.space_after = Pt(size * 0.15)
        return box

    def _filled_rect(slide, left, top, width, height, fill_color, border=False):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if not border:
            shape.line.fill.background()
        return shape

    def _add_cropped_picture(slide, image_path, left, top, box_width, box_height):
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        box_w_emu = int(box_width)
        box_h_emu = int(box_height)
        scale_w = box_w_emu / img_w
        scale_h = box_h_emu / img_h
        if scale_w > scale_h:
            scaled_w = box_w_emu
            scaled_h = int(img_h * scale_w)
            pic = slide.shapes.add_picture(str(image_path), left, top, scaled_w, scaled_h)
            overflow = scaled_h - box_h_emu
            crop_each = overflow / scaled_h / 2
            pic.crop_top = crop_each
            pic.crop_bottom = crop_each
        else:
            scaled_w = int(img_w * scale_h)
            scaled_h = box_h_emu
            pic = slide.shapes.add_picture(str(image_path), left, top, scaled_w, scaled_h)
            overflow = scaled_w - box_w_emu
            crop_each = overflow / scaled_w / 2
            pic.crop_left = crop_each
            pic.crop_right = crop_each
        pic.left = int(left)
        pic.top = int(top)
        pic.width = int(box_w_emu)
        pic.height = int(box_h_emu)
        return pic

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    margin = Inches(0.4)
    content_width = SLIDE_WIDTH - 2 * margin
    y = Inches(0.3)

    photos = [Path(p) for p in (photo_paths or []) if Path(p).exists()]

    # Header: Logo + headline
    if branding.has_logo:
        slide.shapes.add_picture(branding.logo_path, margin, y, Inches(1.6), Inches(0.6))

    _textbox(slide, Inches(2.3), y, Inches(5.8), Inches(0.65),
             branding.headline, size=34, bold=True, italic=True,
             color=ACCENT, align=PP_ALIGN.RIGHT, font="Georgia")

    y += Inches(0.75)

    # Hero photo
    hero_height = Inches(3.5)
    if photos:
        _add_cropped_picture(slide, photos[0], margin, y, content_width, hero_height)
    else:
        _filled_rect(slide, margin, y, content_width, hero_height, LIGHT_GREY)
        _textbox(slide, margin, y + Inches(1.4), content_width, Inches(0.6),
                 "[ HERO PHOTO ]", size=20, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Address overlay
    address = listing_data.get("address_full", "")
    _filled_rect(slide, margin, y, content_width, Inches(0.4), DARK_GREY)
    _textbox(slide, margin + Inches(0.15), y + Inches(0.02),
             content_width - Inches(0.3), Inches(0.36),
             address, size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    y += hero_height + Inches(0.1)

    # Secondary photo grid
    if len(photos) >= 2:
        grid_height = Inches(1.6)
        grid_gap = Inches(0.1)
        grid_width = (content_width - grid_gap) / 2
        _add_cropped_picture(slide, photos[1], margin, y, int(grid_width), grid_height)
        if len(photos) >= 3:
            _add_cropped_picture(slide, photos[2], margin + int(grid_width) + grid_gap, y,
                                 int(grid_width), grid_height)
        else:
            _filled_rect(slide, margin + int(grid_width) + grid_gap, y,
                         int(grid_width), grid_height, LIGHT_GREY)
        y += grid_height + Inches(0.1)

    # Price / specs bar
    specs = _build_specs(listing_data)
    _filled_rect(slide, margin, y, content_width, Inches(0.4), DARK_GREY)
    _textbox(slide, margin, y + Inches(0.02), content_width, Inches(0.36),
             "   |   ".join(specs), size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    y += Inches(0.5)

    # Feature checklist
    features = (listing_data.get("features") or [])[:8]
    if features:
        mid = (len(features) + 1) // 2
        col1 = features[:mid]
        col2 = features[mid:]
        col_w = content_width // 2
        row_h = Inches(0.22)
        _multiline_textbox(slide, margin + Inches(0.15), y,
                           col_w - Inches(0.15), row_h * len(col1),
                           [f"\u25A0  {f}" for f in col1], size=10, color=DARK_GREY, bold=True)
        _multiline_textbox(slide, margin + col_w + Inches(0.15), y,
                           col_w - Inches(0.15), row_h * len(col2),
                           [f"\u25A0  {f}" for f in col2], size=10, color=DARK_GREY, bold=True)
        y += row_h * max(len(col1), len(col2)) + Inches(0.1)

    # Body copy
    body = _extract_body_copy(flyer_text)
    has_photos = len(photos) >= 2
    space_below = Inches(1.0) + Inches(0.5) + Inches(0.2)
    body_box_h = SLIDE_HEIGHT - y - space_below - Inches(0.35)
    body_box_h = max(Inches(0.5), min(body_box_h, Inches(1.2)))
    body_font = 10 if has_photos else 11
    body_line_spacing = Pt(14) if has_photos else Pt(16)

    _filled_rect(slide, margin + Inches(0.6), y,
                 content_width - Inches(1.2), Emu(6350), LIGHT_GREY)
    y += Inches(0.12)

    copy_box = _textbox(slide, margin + Inches(0.4), y,
                        content_width - Inches(0.8), body_box_h,
                        body, size=body_font, color=DARK_GREY,
                        align=PP_ALIGN.CENTER, font="Georgia", italic=True)
    copy_box.text_frame.word_wrap = True
    for para in copy_box.text_frame.paragraphs:
        para.line_spacing = body_line_spacing
    y += body_box_h + Inches(0.03)

    _filled_rect(slide, margin + Inches(0.6), y,
                 content_width - Inches(1.2), Emu(6350), LIGHT_GREY)
    y += Inches(0.15)

    # Bottom section: QR | Agent | Logo
    agent_name = listing_data.get("listing_agent_name", "")
    agent_email = listing_data.get("listing_agent_email", "")
    agent_phone = listing_data.get("listing_agent_phone", "")

    qr_url = _build_qr_url(listing_data, branding)
    qr_path = _generate_qr(qr_url)
    qr_size = Inches(0.95)
    slide.shapes.add_picture(qr_path, margin + Inches(0.1), y, qr_size, qr_size)
    _textbox(slide, margin, y + qr_size + Inches(0.02), Inches(1.2), Inches(0.18),
             "Scan for Details", size=7, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

    agent_left = margin + Inches(1.5)
    if agent_name:
        _textbox(slide, agent_left, y, Inches(3.5), Inches(0.3),
                 agent_name, size=15, bold=True, italic=True, color=BLACK, font="Georgia")
    _textbox(slide, agent_left, y + Inches(0.3), Inches(3.5), Inches(0.2),
             branding.brokerage_name, size=10, color=MID_GREY)
    contact_parts = [p for p in [agent_phone, agent_email] if p]
    if contact_parts:
        _textbox(slide, agent_left, y + Inches(0.52), Inches(3.5), Inches(0.2),
                 "  |  ".join(contact_parts), size=10, bold=True, color=DARK_GREY)

    if branding.has_logo:
        slide.shapes.add_picture(
            branding.logo_path, SLIDE_WIDTH - margin - Inches(1.7), y + Inches(0.05),
            Inches(1.6), Inches(0.6),
        )

    y += Inches(1.0)

    # Footer
    _filled_rect(slide, margin, y, content_width, Emu(12700), BLACK)
    y += Inches(0.06)
    col_w = content_width // 3

    _textbox(slide, margin, y, col_w, Inches(0.45),
             branding.tagline or branding.brokerage_name,
             size=8, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    _textbox(slide, margin + col_w, y, col_w, Inches(0.45),
             branding.brokerage_address or branding.brokerage_name,
             size=8, color=BLACK, align=PP_ALIGN.CENTER)
    footer_right = "\n".join(filter(None, [branding.brokerage_phone, branding.brokerage_website]))
    _textbox(slide, margin + 2 * col_w, y, col_w, Inches(0.45),
             footer_right, size=8, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)

    # Page border
    border = slide.shapes.add_shape(1, Inches(0.2), Inches(0.2),
                                    SLIDE_WIDTH - Inches(0.4), SLIDE_HEIGHT - Inches(0.4))
    border.fill.background()
    border.line.color.rgb = BLACK
    border.line.width = Pt(1.5)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── PDF Builder ───────────────────────────────────────────────────────

def build_flyer_pdf(
    listing_data: dict,
    flyer_text: str,
    branding: BrandingConfig,
    photo_paths: list[Path] | None = None,
) -> io.BytesIO:
    """Generate a branded PDF flyer. Returns bytes buffer.

    Args:
        listing_data: Listing data dict (or Listing model attributes).
        flyer_text: AI-generated flyer content.
        branding: Visual branding configuration.
        photo_paths: Local image file paths (hero + secondary).

    Returns:
        BytesIO buffer containing the PDF file.
    """
    from fpdf import FPDF

    BLACK = (0x00, 0x00, 0x00)
    WHITE = (0xFF, 0xFF, 0xFF)
    DARK_GREY = (0x33, 0x33, 0x33)
    MID_GREY = (0x66, 0x66, 0x66)
    LIGHT_GREY = (0xE8, 0xE8, 0xE8)
    ACCENT = branding.accent_rgb

    PAGE_W = 215.9   # mm
    PAGE_H = 279.4   # mm
    MARGIN = 10.16   # mm
    CONTENT_W = PAGE_W - 2 * MARGIN

    def _in(inches: float) -> float:
        return inches * 25.4

    def _fit_image_in_box(pdf, image_path, x, y, box_w, box_h):
        with Image.open(str(image_path)) as img:
            img_w_px, img_h_px = img.size
        scale = max(box_w / img_w_px, box_h / img_h_px)
        render_w = img_w_px * scale
        render_h = img_h_px * scale
        offset_x = x - (render_w - box_w) / 2
        offset_y = y - (render_h - box_h) / 2
        with pdf.rect_clip(x=x, y=y, w=box_w, h=box_h):
            pdf.image(str(image_path), x=offset_x, y=offset_y, w=render_w, h=render_h)

    pdf = FPDF(orientation="P", unit="mm", format="Letter")
    pdf.set_auto_page_break(auto=False)
    pdf.add_page()

    photos = [Path(p) for p in (photo_paths or []) if Path(p).exists()]
    y = _in(0.3)

    # Page border
    pdf.set_draw_color(*BLACK)
    pdf.set_line_width(0.5)
    pdf.rect(_in(0.2), _in(0.2), PAGE_W - _in(0.4), PAGE_H - _in(0.4))

    # Header
    if branding.has_logo:
        pdf.image(branding.logo_path, x=MARGIN, y=y, w=_in(1.6), h=_in(0.6))

    pdf.set_font("Times", "BI", 28)
    pdf.set_text_color(*ACCENT)
    text_w = pdf.get_string_width(branding.headline)
    pdf.text(x=MARGIN + CONTENT_W - text_w, y=y + _in(0.42), text=branding.headline)

    y += _in(0.75)

    # Hero photo
    hero_h = _in(3.5)
    if photos:
        _fit_image_in_box(pdf, photos[0], MARGIN, y, CONTENT_W, hero_h)
    else:
        pdf.set_fill_color(*LIGHT_GREY)
        pdf.set_draw_color(*MID_GREY)
        pdf.set_line_width(0.15)
        pdf.rect(MARGIN, y, CONTENT_W, hero_h, style="DF")
        pdf.set_font("Helvetica", "", 16)
        pdf.set_text_color(*MID_GREY)
        placeholder = "[ HERO PHOTO ]"
        pw = pdf.get_string_width(placeholder)
        pdf.text(x=MARGIN + (CONTENT_W - pw) / 2, y=y + hero_h / 2 + 2, text=placeholder)

    # Address overlay
    address = _sanitize_text(listing_data.get("address_full", ""))
    pdf.set_fill_color(*DARK_GREY)
    addr_bar_h = _in(0.4)
    pdf.rect(MARGIN, y, CONTENT_W, addr_bar_h, style="F")
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_text_color(*WHITE)
    pdf.text(x=MARGIN + _in(0.15), y=y + addr_bar_h - _in(0.1), text=address)

    y += hero_h + _in(0.1)

    # Secondary photos
    if len(photos) >= 2:
        grid_h = _in(1.6)
        grid_gap = _in(0.1)
        grid_w = (CONTENT_W - grid_gap) / 2
        _fit_image_in_box(pdf, photos[1], MARGIN, y, grid_w, grid_h)
        if len(photos) >= 3:
            _fit_image_in_box(pdf, photos[2], MARGIN + grid_w + grid_gap, y, grid_w, grid_h)
        else:
            pdf.set_fill_color(*LIGHT_GREY)
            pdf.rect(MARGIN + grid_w + grid_gap, y, grid_w, grid_h, style="F")
        y += grid_h + _in(0.1)

    # Specs bar
    specs = _build_specs(listing_data)
    specs_bar_h = _in(0.4)
    pdf.set_fill_color(*DARK_GREY)
    pdf.rect(MARGIN, y, CONTENT_W, specs_bar_h, style="F")
    specs_text = _sanitize_text("   |   ".join(specs))
    pdf.set_font("Helvetica", "B", 12)
    pdf.set_text_color(*WHITE)
    sw = pdf.get_string_width(specs_text)
    pdf.text(x=MARGIN + (CONTENT_W - sw) / 2, y=y + specs_bar_h - _in(0.1), text=specs_text)
    y += _in(0.5)

    # Feature checklist
    features = (listing_data.get("features") or [])[:8]
    if features:
        mid = (len(features) + 1) // 2
        col1 = features[:mid]
        col2 = features[mid:]
        col_w = CONTENT_W / 2
        row_h = _in(0.22)
        pdf.set_font("Helvetica", "B", 8.5)
        pdf.set_text_color(*DARK_GREY)
        for i, feat in enumerate(col1):
            pdf.text(x=MARGIN + _in(0.15), y=y + row_h * i + row_h * 0.7,
                     text=_sanitize_text(f"\u00bb  {feat}"))
        for i, feat in enumerate(col2):
            pdf.text(x=MARGIN + col_w + _in(0.15), y=y + row_h * i + row_h * 0.7,
                     text=_sanitize_text(f"\u00bb  {feat}"))
        y += row_h * max(len(col1), len(col2)) + _in(0.1)

    # Body copy
    body = _extract_body_copy(flyer_text)
    pdf.set_fill_color(*LIGHT_GREY)
    pdf.rect(MARGIN + _in(0.6), y, CONTENT_W - _in(1.2), 0.5, style="F")
    y += _in(0.15)

    pdf.set_font("Times", "I", 10)
    pdf.set_text_color(*DARK_GREY)
    pdf.set_xy(MARGIN + _in(0.4), y)
    pdf.multi_cell(w=CONTENT_W - _in(0.8), h=5.5, text=_sanitize_text(body), align="C")
    y = pdf.get_y() + _in(0.08)

    pdf.set_fill_color(*LIGHT_GREY)
    pdf.rect(MARGIN + _in(0.6), y, CONTENT_W - _in(1.2), 0.5, style="F")
    y += _in(0.2)

    # Bottom: QR | Agent | Logo
    agent_name = _sanitize_text(listing_data.get("listing_agent_name", ""))
    agent_email = _sanitize_text(listing_data.get("listing_agent_email", ""))
    agent_phone = _sanitize_text(listing_data.get("listing_agent_phone", ""))

    qr_url = _build_qr_url(listing_data, branding)
    qr_path = _generate_qr(qr_url)
    qr_size = _in(0.95)
    pdf.image(qr_path, x=MARGIN + _in(0.1), y=y, w=qr_size, h=qr_size)

    pdf.set_font("Helvetica", "B", 6)
    pdf.set_text_color(*ACCENT)
    scan_text = "Scan for Details"
    scan_w = pdf.get_string_width(scan_text)
    pdf.text(x=MARGIN + _in(0.1) + qr_size / 2 - scan_w / 2,
             y=y + qr_size + _in(0.12), text=scan_text)

    agent_x = MARGIN + _in(1.5)
    if agent_name:
        pdf.set_font("Times", "BI", 13)
        pdf.set_text_color(*BLACK)
        pdf.text(x=agent_x, y=y + _in(0.2), text=agent_name)

    pdf.set_font("Helvetica", "", 8.5)
    pdf.set_text_color(*MID_GREY)
    pdf.text(x=agent_x, y=y + _in(0.42), text=_sanitize_text(branding.brokerage_name))

    contact_parts = [p for p in [agent_phone, agent_email] if p]
    if contact_parts:
        pdf.set_font("Helvetica", "B", 8.5)
        pdf.set_text_color(*DARK_GREY)
        pdf.text(x=agent_x, y=y + _in(0.6), text="  |  ".join(contact_parts))

    if branding.has_logo:
        pdf.image(branding.logo_path, x=PAGE_W - MARGIN - _in(1.6),
                  y=y + _in(0.05), w=_in(1.6), h=_in(0.6))

    y += _in(1.15)

    # Footer
    pdf.set_fill_color(*BLACK)
    pdf.rect(MARGIN, y, CONTENT_W, 1.0, style="F")
    y += _in(0.06)

    col_w = CONTENT_W / 3
    fs = 7

    pdf.set_font("Helvetica", "B", fs)
    pdf.set_text_color(*BLACK)
    tagline = _sanitize_text(branding.tagline or branding.brokerage_name)
    pdf.text(x=MARGIN, y=y + 3.5, text=tagline)

    if branding.brokerage_address:
        pdf.set_font("Helvetica", "", fs)
        for i, line in enumerate(branding.brokerage_address.split("\n")[:3]):
            lw = pdf.get_string_width(line)
            pdf.text(x=MARGIN + col_w + (col_w - lw) / 2,
                     y=y + 3.0 + i * (fs * 0.42), text=_sanitize_text(line))

    pdf.set_font("Helvetica", "B", fs)
    right_lines = [l for l in [branding.brokerage_phone, branding.brokerage_website] if l]
    for i, line in enumerate(right_lines):
        lw = pdf.get_string_width(line)
        pdf.text(x=MARGIN + 2 * col_w + (col_w - lw),
                 y=y + 3.5 + i * (fs * 0.42), text=_sanitize_text(line))

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


# ── Convenience ───────────────────────────────────────────────────────

class FlyerService:
    """High-level flyer generation service."""

    def __init__(self, branding: BrandingConfig):
        self.branding = branding

    def generate_pptx(
        self,
        listing_data: dict,
        flyer_text: str,
        photo_paths: list[Path] | None = None,
    ) -> io.BytesIO:
        return build_flyer_pptx(listing_data, flyer_text, self.branding, photo_paths)

    def generate_pdf(
        self,
        listing_data: dict,
        flyer_text: str,
        photo_paths: list[Path] | None = None,
    ) -> io.BytesIO:
        return build_flyer_pdf(listing_data, flyer_text, self.branding, photo_paths)
